"""Tests for theme functionality."""

import io
import os
import tempfile

import pytest
from pptx import Presentation

from app.schemas.theme import (
    DEFAULT_PRIMARY_COLOR,
    MAX_FOOTER_LENGTH,
    MAX_LOGO_SIZE_BYTES,
    Theme,
    ThemeUpdate,
    get_default_theme,
    merge_theme_with_defaults,
)
from app.services.pptx_generator import PPTXGenerator, hex_to_rgb


class TestThemeSchema:
    """Tests for Theme Pydantic schema."""

    def test_default_theme(self):
        """Test default theme values."""
        theme = get_default_theme()
        # primary_color is normalized to uppercase by validator
        assert theme.primary_color.upper() == DEFAULT_PRIMARY_COLOR.upper()
        assert theme.footer_text is None
        assert theme.logo_uri is None
        assert theme.show_logo is True
        assert theme.show_footer is True

    def test_valid_hex_colors(self):
        """Test valid hex color formats."""
        valid_colors = ["#667EEA", "#FF0000", "#00ff00", "#123456", "#ABCDEF"]
        for color in valid_colors:
            theme = Theme(primary_color=color)
            assert theme.primary_color == color.upper()

    def test_invalid_hex_colors(self):
        """Test invalid hex color formats return 422."""
        invalid_colors = [
            "667EEA",  # Missing #
            "#FFF",  # Too short (3 digits)
            "#GGGGGG",  # Invalid hex chars
            "red",  # Named color
            "#12345",  # 5 digits
            "#1234567",  # 7 digits
        ]
        for color in invalid_colors:
            with pytest.raises(ValueError) as exc_info:
                Theme(primary_color=color)
            assert "Invalid color format" in str(exc_info.value)

    def test_footer_text_max_length(self):
        """Test footer text length validation."""
        # Valid: exactly at max length
        valid_footer = "a" * MAX_FOOTER_LENGTH
        theme = Theme(footer_text=valid_footer)
        assert len(theme.footer_text) == MAX_FOOTER_LENGTH

        # Invalid: exceeds max length
        invalid_footer = "a" * (MAX_FOOTER_LENGTH + 1)
        with pytest.raises(ValueError):
            Theme(footer_text=invalid_footer)

    def test_footer_text_empty_string_becomes_none(self):
        """Test that empty footer text is converted to None."""
        theme = Theme(footer_text="")
        assert theme.footer_text is None

        theme = Theme(footer_text="   ")  # Whitespace only
        assert theme.footer_text is None

    def test_merge_theme_with_defaults_none(self):
        """Test merging None theme_json returns default."""
        theme = merge_theme_with_defaults(None)
        default = get_default_theme()
        assert theme.primary_color == default.primary_color
        assert theme.show_logo == default.show_logo

    def test_merge_theme_with_defaults_partial(self):
        """Test merging partial theme_json applies defaults."""
        partial = {"primary_color": "#FF0000"}
        theme = merge_theme_with_defaults(partial)
        assert theme.primary_color == "#FF0000"
        assert theme.show_logo is True  # Default
        assert theme.show_footer is True  # Default

    def test_merge_theme_with_defaults_full(self):
        """Test merging full theme_json preserves all values."""
        full = {
            "primary_color": "#123456",
            "footer_text": "Test Footer",
            "logo_uri": "s3://bucket/logo.png",
            "show_logo": False,
            "show_footer": False,
        }
        theme = merge_theme_with_defaults(full)
        assert theme.primary_color == "#123456"
        assert theme.footer_text == "Test Footer"
        assert theme.logo_uri == "s3://bucket/logo.png"
        assert theme.show_logo is False
        assert theme.show_footer is False


class TestThemeUpdate:
    """Tests for ThemeUpdate schema."""

    def test_partial_update(self):
        """Test that ThemeUpdate allows partial updates."""
        update = ThemeUpdate(primary_color="#FF0000")
        assert update.primary_color == "#FF0000"
        assert update.footer_text is None
        assert update.show_logo is None

    def test_invalid_color_in_update(self):
        """Test invalid color in update raises error."""
        with pytest.raises(ValueError):
            ThemeUpdate(primary_color="invalid")


class TestHexToRgb:
    """Tests for hex_to_rgb conversion."""

    def test_hex_to_rgb_basic(self):
        """Test basic hex to RGB conversion."""
        rgb = hex_to_rgb("#FF0000")  # Red
        # RGBColor is tuple-like: (r, g, b)
        assert rgb[0] == 255  # red
        assert rgb[1] == 0  # green
        assert rgb[2] == 0  # blue

    def test_hex_to_rgb_with_hash(self):
        """Test hex to RGB with hash prefix."""
        rgb = hex_to_rgb("#00FF00")  # Green
        assert rgb[0] == 0  # red
        assert rgb[1] == 255  # green
        assert rgb[2] == 0  # blue

    def test_hex_to_rgb_mixed(self):
        """Test hex to RGB with mixed values."""
        rgb = hex_to_rgb("#667EEA")  # Default purple
        assert rgb[0] == 102  # red
        assert rgb[1] == 126  # green
        assert rgb[2] == 234  # blue


class TestPPTXGeneratorWithTheme:
    """Tests for PPTX generation with theme support."""

    def test_generate_with_default_theme(self):
        """Test PPTX generation with default theme."""
        generator = PPTXGenerator()
        steps_data = {
            "title": "Test Manual",
            "goal": "Test Goal",
            "steps": [
                {
                    "no": 1,
                    "telop": "Step 1",
                    "action": "Do something",
                    "narration": "Narration text",
                    "start": "00:00",
                    "end": "00:10",
                    "frame_file": "step_001.png",
                }
            ],
        }

        pptx_bytes = generator.generate(steps_data, {})
        assert len(pptx_bytes) > 0

        # Verify it's a valid PPTX
        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) >= 2  # Title + at least 1 step

    def test_generate_with_custom_primary_color(self):
        """Test PPTX generation applies custom primary color to titles."""
        generator = PPTXGenerator()
        theme = Theme(primary_color="#FF0000")  # Red
        steps_data = {
            "title": "Test Manual",
            "steps": [
                {
                    "no": 1,
                    "telop": "Step 1",
                    "action": "Do something",
                    "start": "00:00",
                    "end": "00:10",
                    "frame_file": "step_001.png",
                }
            ],
        }

        pptx_bytes = generator.generate(steps_data, {}, theme=theme)

        # Verify the PPTX contains red color
        prs = Presentation(io.BytesIO(pptx_bytes))

        # Find title text box in first slide and check color
        title_slide = prs.slides[0]
        found_red = False
        for shape in title_slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font.color.rgb:
                        # RGBColor is tuple-like: (r, g, b)
                        if para.font.color.rgb[0] == 255:
                            found_red = True
                            break

        assert found_red, "Primary color should be applied to title"

    def test_generate_with_footer_text(self):
        """Test PPTX generation includes footer text when enabled."""
        generator = PPTXGenerator()
        theme = Theme(footer_text="© 2026 Test Company", show_footer=True)
        steps_data = {
            "title": "Test Manual",
            "steps": [],
        }

        pptx_bytes = generator.generate(steps_data, {}, theme=theme)
        prs = Presentation(io.BytesIO(pptx_bytes))

        # Check that footer text is in at least one slide
        footer_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if "© 2026 Test Company" in para.text:
                            footer_found = True
                            break

        assert footer_found, "Footer text should be included in PPTX"

    def test_generate_without_footer_when_disabled(self):
        """Test PPTX generation excludes footer when show_footer=False."""
        generator = PPTXGenerator()
        theme = Theme(footer_text="© 2026 Test Company", show_footer=False)
        steps_data = {
            "title": "Test Manual",
            "steps": [],
        }

        pptx_bytes = generator.generate(steps_data, {}, theme=theme)
        prs = Presentation(io.BytesIO(pptx_bytes))

        # Check that footer text is NOT in any slide
        footer_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if "© 2026 Test Company" in para.text:
                            footer_found = True

        assert not footer_found, "Footer should not be included when disabled"

    def test_generate_with_logo(self):
        """Test PPTX generation includes logo when provided."""
        generator = PPTXGenerator()
        theme = Theme(show_logo=True, logo_uri="s3://bucket/logo.png")
        steps_data = {
            "title": "Test Manual",
            "steps": [],
        }

        # Create a temp PNG file for logo
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            # Create a minimal valid PNG
            import struct
            import zlib

            def create_minimal_png():
                """Create a minimal valid 1x1 PNG."""
                signature = b"\x89PNG\r\n\x1a\n"
                # IHDR chunk
                width = 1
                height = 1
                bit_depth = 8
                color_type = 2  # RGB
                ihdr_data = struct.pack(">IIBBBBB", width, height, bit_depth, color_type, 0, 0, 0)
                ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
                ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
                # IDAT chunk (1x1 red pixel)
                raw_data = b"\x00\xff\x00\x00"  # filter byte + RGB
                compressed = zlib.compress(raw_data)
                idat_crc = zlib.crc32(b"IDAT" + compressed) & 0xFFFFFFFF
                idat = (
                    struct.pack(">I", len(compressed))
                    + b"IDAT"
                    + compressed
                    + struct.pack(">I", idat_crc)
                )
                # IEND chunk
                iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
                iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
                return signature + ihdr + idat + iend

            f.write(create_minimal_png())
            logo_path = f.name

        try:
            pptx_bytes = generator.generate(steps_data, {}, theme=theme, logo_path=logo_path)
            prs = Presentation(io.BytesIO(pptx_bytes))

            # Check that there's at least one picture shape
            picture_found = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        picture_found = True
                        break

            assert picture_found, "Logo should be included in PPTX"
        finally:
            os.unlink(logo_path)

    def test_generate_without_logo_when_disabled(self):
        """Test PPTX generation excludes logo when show_logo=False."""
        generator = PPTXGenerator()
        theme = Theme(show_logo=False, logo_uri="s3://bucket/logo.png")
        steps_data = {
            "title": "Test Manual",
            "steps": [],
        }

        # Create a temp PNG file
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            import struct
            import zlib

            def create_minimal_png():
                signature = b"\x89PNG\r\n\x1a\n"
                width = 1
                height = 1
                bit_depth = 8
                color_type = 2
                ihdr_data = struct.pack(">IIBBBBB", width, height, bit_depth, color_type, 0, 0, 0)
                ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
                ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
                raw_data = b"\x00\xff\x00\x00"
                compressed = zlib.compress(raw_data)
                idat_crc = zlib.crc32(b"IDAT" + compressed) & 0xFFFFFFFF
                idat = (
                    struct.pack(">I", len(compressed))
                    + b"IDAT"
                    + compressed
                    + struct.pack(">I", idat_crc)
                )
                iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
                iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
                return signature + ihdr + idat + iend

            f.write(create_minimal_png())
            logo_path = f.name

        try:
            pptx_bytes = generator.generate(steps_data, {}, theme=theme, logo_path=logo_path)
            prs = Presentation(io.BytesIO(pptx_bytes))

            # Check that there are no picture shapes (logo should be excluded)
            picture_found = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        picture_found = True

            assert not picture_found, "Logo should not be included when disabled"
        finally:
            os.unlink(logo_path)


class TestThemeAPIValidation:
    """Tests for Theme API validation logic using conftest fixtures."""

    def test_put_theme_validates_invalid_color_format(self, client):
        """Test PUT /api/jobs/{job_id}/theme rejects invalid color format."""
        import uuid

        job_id = uuid.uuid4()
        response = client.put(f"/api/jobs/{job_id}/theme", json={"primary_color": "invalid"})
        assert response.status_code == 422

    def test_put_theme_validates_invalid_color_without_hash(self, client):
        """Test that color without # prefix is rejected."""
        import uuid

        job_id = uuid.uuid4()
        response = client.put(f"/api/jobs/{job_id}/theme", json={"primary_color": "FF0000"})
        assert response.status_code == 422

    def test_put_theme_validates_color_too_short(self, client):
        """Test that color with too few digits is rejected."""
        import uuid

        job_id = uuid.uuid4()
        response = client.put(f"/api/jobs/{job_id}/theme", json={"primary_color": "#FFF"})
        assert response.status_code == 422

    def test_get_theme_returns_default_for_succeeded_job(self, client, succeeded_job):
        """Test GET /api/jobs/{job_id}/theme returns default when no theme set."""
        response = client.get(f"/api/jobs/{succeeded_job.id}/theme")
        assert response.status_code == 200
        data = response.json()
        assert data["primary_color"].upper() == DEFAULT_PRIMARY_COLOR.upper()
        assert data["show_logo"] is True
        assert data["show_footer"] is True

    def test_put_theme_updates_for_succeeded_job(self, client, succeeded_job):
        """Test PUT /api/jobs/{job_id}/theme updates theme successfully."""
        response = client.put(
            f"/api/jobs/{succeeded_job.id}/theme",
            json={"primary_color": "#FF0000", "footer_text": "Test Footer"},
        )
        assert response.status_code == 200
        data = response.json()
        assert data["theme"]["primary_color"] == "#FF0000"
        assert data["theme"]["footer_text"] == "Test Footer"

    def test_logo_upload_validates_file_type_txt(self, client, succeeded_job):
        """Test POST /api/jobs/{job_id}/theme/logo rejects non-image files."""
        response = client.post(
            f"/api/jobs/{succeeded_job.id}/theme/logo",
            files={"logo_file": ("test.txt", b"not an image", "text/plain")},
        )
        assert response.status_code == 400
        assert "Unsupported" in response.json()["detail"]

    def test_logo_upload_validates_file_size_limit(self, client, succeeded_job):
        """Test POST /api/jobs/{job_id}/theme/logo validates file size."""
        large_file = b"x" * (MAX_LOGO_SIZE_BYTES + 1)
        response = client.post(
            f"/api/jobs/{succeeded_job.id}/theme/logo",
            files={"logo_file": ("test.png", large_file, "image/png")},
        )
        assert response.status_code == 400
        assert "too large" in response.json()["detail"]

    def test_logo_upload_succeeds_with_valid_file(self, client, succeeded_job):
        """Test POST /api/jobs/{job_id}/theme/logo succeeds with valid PNG."""
        import struct
        import zlib

        def create_minimal_png():
            signature = b"\x89PNG\r\n\x1a\n"
            width = 1
            height = 1
            bit_depth = 8
            color_type = 2
            ihdr_data = struct.pack(">IIBBBBB", width, height, bit_depth, color_type, 0, 0, 0)
            ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
            ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
            raw_data = b"\x00\xff\x00\x00"
            compressed = zlib.compress(raw_data)
            idat_crc = zlib.crc32(b"IDAT" + compressed) & 0xFFFFFFFF
            idat = (
                struct.pack(">I", len(compressed))
                + b"IDAT"
                + compressed
                + struct.pack(">I", idat_crc)
            )
            iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
            iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
            return signature + ihdr + idat + iend

        png_data = create_minimal_png()
        response = client.post(
            f"/api/jobs/{succeeded_job.id}/theme/logo",
            files={"logo_file": ("logo.png", png_data, "image/png")},
        )
        assert response.status_code == 200
        assert "logo_uri" in response.json()
