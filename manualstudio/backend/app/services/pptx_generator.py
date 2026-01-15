"""PPTX generation service."""

import os
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.logging import get_logger
from app.schemas.theme import Theme, get_default_theme

logger = get_logger(__name__)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor.

    Args:
        hex_color: Hex color string (e.g., "#667EEA")

    Returns:
        RGBColor object
    """
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


class PPTXGenerator:
    """Generate PPTX from steps.json."""

    def __init__(self):
        self.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.slide_height = Inches(7.5)

    def generate(
        self,
        steps_data: dict,
        frame_paths: dict[str, str],
        output_path: str | None = None,
        theme: Theme | None = None,
        logo_path: str | None = None,
    ) -> bytes:
        """
        Generate PPTX from steps data.

        Args:
            steps_data: steps.json data
            frame_paths: Dict mapping frame_file to local file path
            output_path: Optional path to save PPTX
            theme: Optional Theme object for customization
            logo_path: Optional local path to logo image file

        Returns:
            PPTX file as bytes
        """
        logger.info("Generating PPTX")

        # Use default theme if not provided
        if theme is None:
            theme = get_default_theme()

        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height

        # Add title slide
        self._add_title_slide(prs, steps_data, theme, logo_path)

        # Add step slides
        for step in steps_data.get("steps", []):
            frame_file = step.get("frame_file", "")
            frame_path = frame_paths.get(frame_file)
            self._add_step_slide(prs, step, frame_path, theme, logo_path)

        # Add common mistakes slide if present
        mistakes = steps_data.get("common_mistakes", [])
        if mistakes:
            self._add_mistakes_slide(prs, mistakes, theme, logo_path)

        # Add quiz slides if present
        quizzes = steps_data.get("quiz", [])
        for quiz in quizzes:
            self._add_quiz_slide(prs, quiz, theme, logo_path)

        # Save to bytes
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        pptx_bytes = buffer.read()

        if output_path:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path, "wb") as f:
                f.write(pptx_bytes)
            logger.info(f"PPTX saved to {output_path}")

        logger.info(f"PPTX generated: {len(steps_data.get('steps', []))} step slides")
        return pptx_bytes

    def _add_logo_to_slide(self, slide, logo_path: str | None, theme: Theme) -> None:
        """Add logo to top-right corner of slide if enabled."""
        if not theme.show_logo or not logo_path or not os.path.exists(logo_path):
            return

        try:
            # Position in top-right corner
            logo_width = Inches(1.2)
            logo_left = self.slide_width - logo_width - Inches(0.3)
            logo_top = Inches(0.2)

            slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)
        except Exception as e:
            logger.warning(f"Failed to add logo to slide: {e}")

    def _add_footer_to_slide(self, slide, theme: Theme) -> None:
        """Add footer text to bottom of slide if enabled."""
        if not theme.show_footer or not theme.footer_text:
            return

        try:
            footer_left = Inches(0.5)
            footer_top = self.slide_height - Inches(0.5)
            footer_width = self.slide_width - Inches(1)
            footer_height = Inches(0.4)

            footer_box = slide.shapes.add_textbox(
                footer_left, footer_top, footer_width, footer_height
            )
            tf = footer_box.text_frame
            p = tf.paragraphs[0]
            p.text = theme.footer_text
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(128, 128, 128)
            p.alignment = PP_ALIGN.CENTER
        except Exception as e:
            logger.warning(f"Failed to add footer to slide: {e}")

    def _add_title_slide(
        self,
        prs: Presentation,
        steps_data: dict,
        theme: Theme,
        logo_path: str | None,
    ) -> None:
        """Add title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        title = steps_data.get("title", "操作マニュアル")
        goal = steps_data.get("goal", "")
        primary_color = hex_to_rgb(theme.primary_color)

        # Title text box
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(12.333)
        height = Inches(1.5)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = primary_color
        p.alignment = PP_ALIGN.CENTER

        # Goal text box
        if goal:
            top = Inches(4.5)
            height = Inches(1)
            goal_box = slide.shapes.add_textbox(left, top, width, height)
            tf = goal_box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = goal
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER

        # Add logo and footer
        self._add_logo_to_slide(slide, logo_path, theme)
        self._add_footer_to_slide(slide, theme)

    def _add_step_slide(
        self,
        prs: Presentation,
        step: dict,
        frame_path: str | None,
        theme: Theme,
        logo_path: str | None,
    ) -> None:
        """Add a step slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        step_no = step.get("no", 0)
        telop = step.get("telop", "")
        action = step.get("action", "")
        target = step.get("target", "")
        caution = step.get("caution", "")
        narration = step.get("narration", "")
        time_range = f"{step.get('start', '00:00')} - {step.get('end', '00:00')}"
        primary_color = hex_to_rgb(theme.primary_color)

        # Title
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(12.333)
        title_height = Inches(0.8)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Step {step_no}: {telop}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = primary_color

        # Image (left side)
        if frame_path and os.path.exists(frame_path):
            img_left = Inches(0.5)
            img_top = Inches(1.3)
            img_width = Inches(7)

            try:
                slide.shapes.add_picture(frame_path, img_left, img_top, width=img_width)
            except Exception as e:
                logger.warning(f"Failed to add image: {e}")

        # Text content (right side)
        text_left = Inches(7.8)
        text_top = Inches(1.3)
        text_width = Inches(5)
        text_height = Inches(5.5)

        text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = text_box.text_frame
        tf.word_wrap = True

        # Time range
        p = tf.paragraphs[0]
        p.text = f"時間: {time_range}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Action
        p = tf.add_paragraph()
        p.text = f"操作: {action}"
        p.font.size = Pt(18)
        p.space_before = Pt(12)

        # Target
        if target and target != "unknown":
            p = tf.add_paragraph()
            p.text = f"対象: {target}"
            p.font.size = Pt(16)
            p.space_before = Pt(8)

        # Caution
        if caution:
            p = tf.add_paragraph()
            p.text = f"注意: {caution}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(200, 50, 50)
            p.space_before = Pt(12)

        # Add narration to speaker notes
        if narration:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = narration

        # Add logo and footer
        self._add_logo_to_slide(slide, logo_path, theme)
        self._add_footer_to_slide(slide, theme)

    def _add_mistakes_slide(
        self,
        prs: Presentation,
        mistakes: list,
        theme: Theme,
        logo_path: str | None,
    ) -> None:
        """Add common mistakes slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        primary_color = hex_to_rgb(theme.primary_color)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "よくあるミスと対処法"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = primary_color

        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(mistakes):
            mistake = item.get("mistake", "")
            fix = item.get("fix", "")

            if i > 0:
                p = tf.add_paragraph()
                p.text = ""
                p.space_before = Pt(16)

            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"ミス: {mistake}"
            p.font.size = Pt(18)
            p.font.bold = True

            p = tf.add_paragraph()
            p.text = f"対処: {fix}"
            p.font.size = Pt(16)
            p.space_before = Pt(4)

        # Add logo and footer
        self._add_logo_to_slide(slide, logo_path, theme)
        self._add_footer_to_slide(slide, theme)

    def _add_quiz_slide(
        self,
        prs: Presentation,
        quiz: dict,
        theme: Theme,
        logo_path: str | None,
    ) -> None:
        """Add quiz slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        primary_color = hex_to_rgb(theme.primary_color)

        quiz_type = quiz.get("type", "text")
        question = quiz.get("q", "")
        answer = quiz.get("a", "")
        choices = quiz.get("choices", [])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "確認クイズ"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = primary_color

        # Question
        q_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.5))
        tf = q_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Q: {question}"
        p.font.size = Pt(24)

        # Choices (for choice type)
        if quiz_type == "choice" and choices:
            c_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(3))
            tf = c_box.text_frame
            tf.word_wrap = True

            for i, choice in enumerate(choices):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{chr(65 + i)}. {choice}"
                p.font.size = Pt(20)
                p.space_before = Pt(8)

        # Answer in speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = f"正解: {answer}"

        # Add logo and footer
        self._add_logo_to_slide(slide, logo_path, theme)
        self._add_footer_to_slide(slide, theme)
